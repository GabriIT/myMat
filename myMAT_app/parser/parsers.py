from __future__ import annotations

import base64
import json
import html
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable
from xml.etree import ElementTree as ET

from langchain_core.documents import Document

from .parser_config import ParseConfig
from .parser_types import FileParseResult, ParseIssue, ParseStatus

try:
    from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader
except Exception:
    Docx2txtLoader = None
    PyPDFLoader = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

try:
    import fitz
except Exception:
    fitz = None

try:
    import numpy as np
except Exception:
    np = None

try:
    from rapidocr_onnxruntime import RapidOCR
except Exception:
    RapidOCR = None

try:
    import requests
except Exception:
    requests = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

PDF_PAGE_DENSITY_MIN = 30.0
PPTX_XML_TEXT_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
PPTX_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PPTX_NOTES_REL_SUFFIX = "/relationships/notesSlide"

OCR_GUIDANCE = (
    "This PDF looks scanned/image-only. Alternatives: run OCR first (ocrmypdf + tesseract), "
    "or use managed OCR (Azure Document Intelligence / Google Document AI)."
)
OCR_LOCAL_GUIDANCE = (
    "Optional local OCR fallback requires `pymupdf`, `rapidocr-onnxruntime`, and `opencv-python`."
)
OFFICE_GUIDANCE = (
    "Try repair/convert via `libreoffice --headless` then parse again, or switch to "
    "`unstructured` parsers for complex Office layouts."
)
LOW_TEXT_GUIDANCE = (
    "Extraction is weak. Re-parse with an alternate backend and keep the better output by "
    "text quality score."
)
PPTX_LOCAL_VISION_GUIDANCE = (
    "Optional PPTX vision enrichment requires `libreoffice`, `poppler-utils`, and `requests`."
)
PPTX_VISION_PROVIDER_GUIDANCE = (
    "PPTX vision enrichment currently supports `ollama` via HTTP at OLLAMA_URL/OLLAMA_HOST."
)


@dataclass(slots=True)
class _PptxSlideRecord:
    slide_number: int
    slide_title: str
    text: str
    content_chars: int
    low_text: bool
    has_notes: bool
    shape_count: int


def _normalize_text(raw: str) -> str:
    text = raw.replace("\r\n", "\n").replace("\r", "\n").replace("\u00a0", " ")
    text = html.unescape(text)
    lines = []
    for line in text.split("\n"):
        compact = re.sub(r"[ \t]+", " ", line).strip()
        lines.append(compact)
    text = "\n".join(lines)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _doc_type_for_path(knowledge_root: Path, path: Path) -> str:
    try:
        relative = path.relative_to(knowledge_root)
    except ValueError:
        return "__external__"
    if len(relative.parts) >= 2:
        return relative.parts[0]
    return "__root__"


def _base_metadata(config: ParseConfig, path: Path) -> dict:
    return {
        "source": str(path),
        "source_name": path.name,
        "source_ext": path.suffix.lower(),
        "doc_type": _doc_type_for_path(config.knowledge_root, path),
    }


def _annotate_docs(
    docs: list[Document],
    parser_used: str,
    fallback_used: bool,
    status: ParseStatus,
    extracted_chars: int,
) -> None:
    for doc in docs:
        doc.metadata["parser_used"] = parser_used
        doc.metadata["fallback_used"] = fallback_used
        doc.metadata["parse_status"] = status
        doc.metadata["extracted_chars"] = extracted_chars


def _result(
    path: Path,
    extension: str,
    status: ParseStatus,
    parser_used: str,
    fallback_used: bool,
    extracted_chars: int,
    documents_count: int,
    issues: list[ParseIssue] | None = None,
) -> FileParseResult:
    return FileParseResult(
        source_path=str(path),
        extension=extension,
        status=status,
        parser_used=parser_used,
        fallback_used=fallback_used,
        extracted_chars=extracted_chars,
        documents_count=documents_count,
        issues=issues or [],
    )


def _iter_supported_files(config: ParseConfig) -> list[Path]:
    if config.include_paths:
        filtered = []
        seen: set[Path] = set()
        for path in config.include_paths:
            resolved = Path(path).expanduser().resolve()
            if resolved in seen:
                continue
            seen.add(resolved)
            if resolved.is_file() and resolved.suffix.lower() in config.supported_extensions:
                filtered.append(resolved)
        return sorted(filtered)

    pattern = "**/*" if config.recursive else "*"
    files = [
        path
        for path in config.knowledge_root.glob(pattern)
        if path.is_file() and path.suffix.lower() in config.supported_extensions
    ]
    return sorted(files)


def _split_pdf_text_pages(text: str) -> list[str]:
    pages = [_normalize_text(chunk) for chunk in text.split("\f")]
    return [page for page in pages if page]


def _parse_pdf_with_pypdf(path: Path, base: dict) -> tuple[list[Document], int, int]:
    if PyPDFLoader is None:
        raise RuntimeError("PyPDFLoader is unavailable (missing langchain_community/pypdf)")
    loader = PyPDFLoader(str(path), mode="page")
    loaded = loader.load()
    docs: list[Document] = []
    for raw in loaded:
        text = _normalize_text(raw.page_content or "")
        if not text:
            continue
        metadata = dict(base)
        metadata.update(raw.metadata or {})
        if "page" in metadata and isinstance(metadata["page"], int):
            metadata["page_number"] = metadata["page"] + 1
        docs.append(Document(page_content=text, metadata=metadata))
    return docs, len(loaded), sum(len(doc.page_content) for doc in docs)


def _parse_pdf_with_pdftotext(path: Path, base: dict) -> tuple[list[Document], int]:
    proc = subprocess.run(
        ["pdftotext", "-layout", str(path), "-"],
        capture_output=True,
        text=True,
        check=False,
    )
    if proc.returncode != 0:
        stderr = proc.stderr.strip() or f"pdftotext exited with {proc.returncode}"
        raise RuntimeError(stderr)

    pages = _split_pdf_text_pages(proc.stdout or "")
    docs: list[Document] = []
    for idx, page_text in enumerate(pages, start=1):
        metadata = dict(base)
        metadata["page_number"] = idx
        docs.append(Document(page_content=page_text, metadata=metadata))
    return docs, sum(len(doc.page_content) for doc in docs)


def _parse_pdf_with_rapidocr(path: Path, base: dict) -> tuple[list[Document], int]:
    if fitz is None or np is None or RapidOCR is None:
        raise RuntimeError("RapidOCR fallback is unavailable")

    ocr_engine = RapidOCR()
    docs: list[Document] = []
    pdf = fitz.open(path)
    try:
        for index, page in enumerate(pdf, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height, pix.width, pix.n
            )
            result, _ = ocr_engine(image)
            if not result:
                continue

            lines: list[str] = []
            for item in result:
                if not isinstance(item, (list, tuple)) or len(item) < 2:
                    continue
                text_part = item[1]
                if isinstance(text_part, (list, tuple)) and text_part:
                    text_part = text_part[0]
                if isinstance(text_part, str):
                    text_clean = _normalize_text(text_part)
                    if text_clean:
                        lines.append(text_clean)

            page_text = _normalize_text("\n".join(lines))
            if not page_text:
                continue

            metadata = dict(base)
            metadata["page_number"] = index
            docs.append(Document(page_content=page_text, metadata=metadata))
    finally:
        pdf.close()

    return docs, sum(len(doc.page_content) for doc in docs)


def parse_pdf_file(path: Path, config: ParseConfig) -> tuple[list[Document], FileParseResult]:
    extension = path.suffix.lower()
    base = _base_metadata(config, path)
    issues: list[ParseIssue] = []

    primary_docs: list[Document] = []
    primary_pages = 0
    primary_chars = 0
    primary_error: Exception | None = None
    try:
        primary_docs, primary_pages, primary_chars = _parse_pdf_with_pypdf(path, base)
    except Exception as exc:
        primary_error = exc
        text = str(exc).lower()
        missing_dep = "pypdf" in text or "no module named" in text
        issue_code = "pdf_primary_parser_unavailable" if missing_dep else "pdf_primary_parser_failed"
        guidance = (
            "Install `pypdf` to enable PyPDFLoader. "
            + LOW_TEXT_GUIDANCE
            if missing_dep
            else LOW_TEXT_GUIDANCE
        )
        issues.append(
            ParseIssue(
                severity="warning",
                code=issue_code,
                message=f"PyPDFLoader failed: {exc}",
                suggested_action=guidance,
            )
        )

    density = (primary_chars / primary_pages) if primary_pages else 0.0
    weak_primary = (
        primary_chars < config.min_chars_pdf
        or (primary_pages > 0 and density < PDF_PAGE_DENSITY_MIN)
    )

    parser_used = "PyPDFLoader"
    fallback_used = False
    chosen_docs = primary_docs
    chosen_chars = primary_chars

    should_try_fallback = primary_error is not None or weak_primary
    if should_try_fallback:
        fallback_used = True
        try:
            fallback_docs, fallback_chars = _parse_pdf_with_pdftotext(path, base)
            if fallback_chars >= chosen_chars:
                chosen_docs = fallback_docs
                chosen_chars = fallback_chars
                parser_used = "pdftotext"
        except Exception as exc:
            issues.append(
                ParseIssue(
                    severity="warning",
                    code="pdf_fallback_failed",
                    message=f"pdftotext fallback failed: {exc}",
                    suggested_action=OCR_GUIDANCE,
                )
            )

    if chosen_chars == 0:
        fallback_used = True
        try:
            ocr_docs, ocr_chars = _parse_pdf_with_rapidocr(path, base)
            if ocr_chars > chosen_chars:
                chosen_docs = ocr_docs
                chosen_chars = ocr_chars
                parser_used = "rapidocr_onnxruntime"
                issues.append(
                    ParseIssue(
                        severity="warning",
                        code="pdf_ocr_applied",
                        message="Applied OCR fallback to recover text from image-based PDF.",
                        suggested_action=(
                            "OCR was used successfully. Verify extracted text quality for this file."
                        ),
                    )
                )
            else:
                issues.append(
                    ParseIssue(
                        severity="warning",
                        code="pdf_ocr_no_text",
                        message="OCR fallback ran but recovered no text.",
                        suggested_action=OCR_GUIDANCE,
                    )
                )
        except Exception:
            issues.append(
                ParseIssue(
                    severity="warning",
                    code="pdf_ocr_unavailable",
                    message="OCR fallback is unavailable in this environment.",
                    suggested_action=f"{OCR_LOCAL_GUIDANCE} {OCR_GUIDANCE}",
                )
            )

    status: ParseStatus = "success"
    if chosen_chars == 0:
        status = "failed"
        issues.append(
            ParseIssue(
                severity="error",
                code="pdf_possible_scan_or_image_only",
                message="No extractable text from PDF.",
                suggested_action=OCR_GUIDANCE,
            )
        )
    else:
        page_count = len(chosen_docs)
        final_density = (chosen_chars / page_count) if page_count else 0.0
        if chosen_chars < config.min_chars_pdf or (
            page_count > 0 and final_density < PDF_PAGE_DENSITY_MIN
        ):
            status = "warning"
            issues.append(
                ParseIssue(
                    severity="warning",
                    code="pdf_possible_scan_or_image_only",
                    message=(
                        "Extracted text is below quality threshold; the document may be mostly "
                        "image-based or hard to parse."
                    ),
                    suggested_action=OCR_GUIDANCE,
                )
            )

    _annotate_docs(chosen_docs, parser_used, fallback_used, status, chosen_chars)
    return chosen_docs, _result(
        path=path,
        extension=extension,
        status=status,
        parser_used=parser_used,
        fallback_used=fallback_used,
        extracted_chars=chosen_chars,
        documents_count=len(chosen_docs),
        issues=issues,
    )


def _parse_docx_with_loader(path: Path, base: dict) -> tuple[list[Document], int]:
    if Docx2txtLoader is None:
        raise RuntimeError("Docx2txtLoader is unavailable (missing docx2txt)")
    loaded = Docx2txtLoader(str(path)).load()
    text = _normalize_text("\n".join(doc.page_content for doc in loaded if doc.page_content))
    if not text:
        return [], 0
    return [Document(page_content=text, metadata=dict(base))], len(text)


def _parse_docx_xml_fallback(path: Path, base: dict) -> tuple[list[Document], int]:
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")
    chunks = re.findall(r"<w:t[^>]*>(.*?)</w:t>", xml)
    text = _normalize_text(" ".join(chunks))
    if not text:
        return [], 0
    return [Document(page_content=text, metadata=dict(base))], len(text)


def parse_docx_file(path: Path, config: ParseConfig) -> tuple[list[Document], FileParseResult]:
    extension = path.suffix.lower()
    base = _base_metadata(config, path)
    issues: list[ParseIssue] = []
    parser_used = "Docx2txtLoader"
    fallback_used = False

    primary_docs: list[Document] = []
    primary_chars = 0
    primary_error: Exception | None = None
    try:
        primary_docs, primary_chars = _parse_docx_with_loader(path, base)
    except Exception as exc:
        primary_error = exc
        text = str(exc).lower()
        missing_dep = "docx2txt" in text or "no module named" in text
        issue_code = "docx_primary_parser_unavailable" if missing_dep else "docx_primary_parser_failed"
        guidance = (
            "Install `docx2txt` for the primary parser path. " + OFFICE_GUIDANCE
            if missing_dep
            else OFFICE_GUIDANCE
        )
        issues.append(
            ParseIssue(
                severity="warning",
                code=issue_code,
                message=f"Docx2txtLoader failed: {exc}",
                suggested_action=guidance,
            )
        )

    chosen_docs = primary_docs
    chosen_chars = primary_chars

    if primary_error is not None or primary_chars < config.min_chars_docx:
        fallback_used = True
        try:
            fallback_docs, fallback_chars = _parse_docx_xml_fallback(path, base)
            if fallback_chars >= chosen_chars:
                chosen_docs = fallback_docs
                chosen_chars = fallback_chars
                parser_used = "docx_xml_fallback"
        except Exception as exc:
            issues.append(
                ParseIssue(
                    severity="warning",
                    code="docx_fallback_failed",
                    message=f"DOCX XML fallback failed: {exc}",
                    suggested_action=OFFICE_GUIDANCE,
                )
            )

    status: ParseStatus = "success"
    if chosen_chars == 0:
        status = "failed"
        issues.append(
            ParseIssue(
                severity="error",
                code="docx_no_extractable_text",
                message="No extractable text from DOCX file.",
                suggested_action=OFFICE_GUIDANCE,
            )
        )
    elif chosen_chars < config.min_chars_docx:
        status = "warning"
        issues.append(
            ParseIssue(
                severity="warning",
                code="docx_low_text",
                message="Extracted DOCX text is below threshold.",
                suggested_action=LOW_TEXT_GUIDANCE,
            )
        )

    _annotate_docs(chosen_docs, parser_used, fallback_used, status, chosen_chars)
    return chosen_docs, _result(
        path=path,
        extension=extension,
        status=status,
        parser_used=parser_used,
        fallback_used=fallback_used,
        extracted_chars=chosen_chars,
        documents_count=len(chosen_docs),
        issues=issues,
    )


def _parse_xlsx_with_openpyxl(path: Path, base: dict) -> tuple[list[Document], int]:
    if load_workbook is None:
        raise RuntimeError("openpyxl is unavailable")
    workbook = load_workbook(path, read_only=True, data_only=True)
    docs: list[Document] = []
    try:
        for sheet in workbook.worksheets:
            row_lines: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = []
                for cell in row:
                    if cell is None:
                        continue
                    text = _normalize_text(str(cell))
                    if text:
                        values.append(text)
                if values:
                    row_lines.append(" | ".join(values))
            if row_lines:
                content = _normalize_text(f"Sheet: {sheet.title}\n" + "\n".join(row_lines))
                if content:
                    metadata = dict(base)
                    metadata["sheet_name"] = sheet.title
                    docs.append(Document(page_content=content, metadata=metadata))
    finally:
        workbook.close()
    return docs, sum(len(doc.page_content) for doc in docs)


def _xlsx_sheet_name_map(zf: zipfile.ZipFile) -> dict[str, str]:
    workbook_xml = "xl/workbook.xml"
    rels_xml = "xl/_rels/workbook.xml.rels"
    if workbook_xml not in zf.namelist() or rels_xml not in zf.namelist():
        return {}

    wb_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    doc_rel_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    wb_root = ET.fromstring(zf.read(workbook_xml))
    rel_root = ET.fromstring(zf.read(rels_xml))

    relation_targets: dict[str, str] = {}
    for rel in rel_root.findall(f"{{{rel_ns}}}Relationship"):
        rid = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        if rid and target:
            normalized = target.lstrip("/")
            if normalized.startswith("../"):
                normalized = normalized[3:]
            if not normalized.startswith("xl/"):
                normalized = f"xl/{normalized}"
            relation_targets[rid] = normalized

    mapping: dict[str, str] = {}
    for sheet in wb_root.findall(f".//{{{wb_ns}}}sheet"):
        name = sheet.attrib.get("name")
        rel_id = sheet.attrib.get(f"{{{doc_rel_ns}}}id")
        if name and rel_id and rel_id in relation_targets:
            mapping[relation_targets[rel_id]] = name
    return mapping


def _shared_strings(zf: zipfile.ZipFile) -> list[str]:
    path = "xl/sharedStrings.xml"
    if path not in zf.namelist():
        return []
    ns = {"s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    root = ET.fromstring(zf.read(path))
    values: list[str] = []
    for si in root.findall("s:si", ns):
        text_parts = []
        for t_node in si.findall(".//s:t", ns):
            text_parts.append(t_node.text or "")
        values.append(_normalize_text("".join(text_parts)))
    return values


def _cell_value(cell: ET.Element, ns: dict[str, str], shared: list[str]) -> str:
    cell_type = cell.attrib.get("t")
    if cell_type == "inlineStr":
        text = "".join(node.text or "" for node in cell.findall(".//s:t", ns))
        return _normalize_text(text)

    value_node = cell.find("s:v", ns)
    if value_node is None or value_node.text is None:
        return ""
    raw = value_node.text.strip()
    if not raw:
        return ""
    if cell_type == "s":
        try:
            idx = int(raw)
            if 0 <= idx < len(shared):
                return _normalize_text(shared[idx])
        except ValueError:
            return ""
    return _normalize_text(raw)


def _parse_xlsx_xml_fallback(path: Path, base: dict) -> tuple[list[Document], int]:
    docs: list[Document] = []
    with zipfile.ZipFile(path) as zf:
        sheets = sorted(
            name
            for name in zf.namelist()
            if name.startswith("xl/worksheets/") and name.endswith(".xml")
        )
        if not sheets:
            return [], 0

        names = _xlsx_sheet_name_map(zf)
        shared = _shared_strings(zf)
        ns = {"s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}

        for sheet_path in sheets:
            root = ET.fromstring(zf.read(sheet_path))
            sheet_name = names.get(sheet_path, Path(sheet_path).stem)
            row_lines: list[str] = []
            for row in root.findall(".//s:sheetData/s:row", ns):
                values = []
                for cell in row.findall("s:c", ns):
                    value = _cell_value(cell, ns, shared)
                    if value:
                        values.append(value)
                if values:
                    row_lines.append(" | ".join(values))
            if row_lines:
                text = _normalize_text(f"Sheet: {sheet_name}\n" + "\n".join(row_lines))
                if text:
                    metadata = dict(base)
                    metadata["sheet_name"] = sheet_name
                    docs.append(Document(page_content=text, metadata=metadata))
    return docs, sum(len(doc.page_content) for doc in docs)


def parse_xlsx_file(path: Path, config: ParseConfig) -> tuple[list[Document], FileParseResult]:
    extension = path.suffix.lower()
    base = _base_metadata(config, path)
    issues: list[ParseIssue] = []
    parser_used = "openpyxl"
    fallback_used = False

    primary_docs: list[Document] = []
    primary_chars = 0
    primary_error: Exception | None = None
    try:
        primary_docs, primary_chars = _parse_xlsx_with_openpyxl(path, base)
    except Exception as exc:
        primary_error = exc
        text = str(exc).lower()
        missing_dep = "openpyxl" in text or "no module named" in text
        issue_code = "xlsx_primary_parser_unavailable" if missing_dep else "xlsx_primary_parser_failed"
        guidance = (
            "Install `openpyxl` for structured workbook parsing. " + OFFICE_GUIDANCE
            if missing_dep
            else OFFICE_GUIDANCE
        )
        issues.append(
            ParseIssue(
                severity="warning",
                code=issue_code,
                message=f"openpyxl parser failed: {exc}",
                suggested_action=guidance,
            )
        )

    chosen_docs = primary_docs
    chosen_chars = primary_chars

    if primary_error is not None or primary_chars < config.min_chars_xlsx:
        fallback_used = True
        try:
            fallback_docs, fallback_chars = _parse_xlsx_xml_fallback(path, base)
            if fallback_chars >= chosen_chars:
                chosen_docs = fallback_docs
                chosen_chars = fallback_chars
                parser_used = "xlsx_xml_fallback"
        except Exception as exc:
            issues.append(
                ParseIssue(
                    severity="warning",
                    code="xlsx_fallback_failed",
                    message=f"XLSX XML fallback failed: {exc}",
                    suggested_action=OFFICE_GUIDANCE,
                )
            )

    status: ParseStatus = "success"
    if chosen_chars == 0:
        status = "failed"
        issues.append(
            ParseIssue(
                severity="error",
                code="xlsx_no_extractable_cells",
                message="No extractable text/cell content from XLSX file.",
                suggested_action=OFFICE_GUIDANCE,
            )
        )
    elif chosen_chars < config.min_chars_xlsx:
        status = "warning"
        issues.append(
            ParseIssue(
                severity="warning",
                code="xlsx_low_text",
                message="Extracted XLSX text is below threshold.",
                suggested_action=LOW_TEXT_GUIDANCE,
            )
        )

    _annotate_docs(chosen_docs, parser_used, fallback_used, status, chosen_chars)
    return chosen_docs, _result(
        path=path,
        extension=extension,
        status=status,
        parser_used=parser_used,
        fallback_used=fallback_used,
        extracted_chars=chosen_chars,
        documents_count=len(chosen_docs),
        issues=issues,
    )


def _pptx_guess_slide_title(slide, default: str) -> str:
    try:
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None:
            title_text = _normalize_text(getattr(title_shape, "text", "") or "")
            if title_text:
                return title_text
    except Exception:
        pass

    for shape in slide.shapes:
        text = _normalize_text(getattr(shape, "text", "") or "")
        if text:
            first_line = text.split("\n", 1)[0].strip()
            if first_line:
                return first_line
    return default


def _pptx_shape_text(shape) -> str:
    sections: list[str] = []
    if getattr(shape, "has_table", False):
        try:
            rows: list[str] = []
            for row in shape.table.rows:
                values = []
                for cell in row.cells:
                    text = _normalize_text(getattr(cell, "text", "") or "")
                    if text:
                        values.append(text)
                if values:
                    rows.append(" | ".join(values))
            if rows:
                sections.append("\n".join(rows))
        except Exception:
            pass

    text_direct = _normalize_text(getattr(shape, "text", "") or "")
    if text_direct:
        sections.append(text_direct)
    return _normalize_text("\n\n".join(sections))


def _pptx_slide_notes(slide) -> str:
    try:
        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide is None:
            return ""
        notes_frame = getattr(notes_slide, "notes_text_frame", None)
        if notes_frame is None:
            return ""
        return _normalize_text(notes_frame.text or "")
    except Exception:
        return ""


def _build_pptx_slide_record(
    *,
    slide_number: int,
    slide_title: str,
    body_text: str,
    notes_text: str,
    shape_count: int,
    min_chars_pptx_slide: int,
) -> _PptxSlideRecord:
    sections = [f"Slide title: {slide_title}"]
    if body_text:
        sections.append(body_text)
    if notes_text:
        sections.append(f"Speaker notes:\n{notes_text}")
    page_text = _normalize_text("\n\n".join(sections))
    content_chars = len(body_text) + len(notes_text)
    return _PptxSlideRecord(
        slide_number=slide_number,
        slide_title=slide_title,
        text=page_text,
        content_chars=content_chars,
        low_text=content_chars < min_chars_pptx_slide,
        has_notes=bool(notes_text),
        shape_count=shape_count,
    )


def _parse_pptx_with_python_pptx(
    path: Path, config: ParseConfig
) -> tuple[list[_PptxSlideRecord], int]:
    if Presentation is None:
        raise RuntimeError("python-pptx is unavailable")

    prs = Presentation(str(path))
    records: list[_PptxSlideRecord] = []
    total_chars = 0

    for idx, slide in enumerate(prs.slides, start=1):
        title = _pptx_guess_slide_title(slide, f"Slide {idx}")
        text_parts: list[str] = []
        for shape in slide.shapes:
            shape_text = _pptx_shape_text(shape)
            if shape_text:
                text_parts.append(shape_text)
        body_text = _normalize_text("\n\n".join(text_parts))
        notes_text = _pptx_slide_notes(slide)
        record = _build_pptx_slide_record(
            slide_number=idx,
            slide_title=title,
            body_text=body_text,
            notes_text=notes_text,
            shape_count=len(slide.shapes),
            min_chars_pptx_slide=config.min_chars_pptx_slide,
        )
        records.append(record)
        total_chars += len(record.text)
    return records, total_chars


def _pptx_slide_xml_paths(zf: zipfile.ZipFile) -> list[tuple[int, str]]:
    pattern = re.compile(r"^ppt/slides/slide(\d+)\.xml$")
    result: list[tuple[int, str]] = []
    for name in zf.namelist():
        match = pattern.match(name)
        if match:
            result.append((int(match.group(1)), name))
    return sorted(result, key=lambda item: item[0])


def _pptx_xml_text(blob: bytes) -> list[str]:
    root = ET.fromstring(blob)
    values = []
    for node in root.findall(".//a:t", PPTX_XML_TEXT_NS):
        text = _normalize_text(node.text or "")
        if text:
            values.append(text)
    return values


def _pptx_notes_xml_path(zf: zipfile.ZipFile, slide_number: int) -> str | None:
    rel_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    if rel_path not in zf.namelist():
        return None
    root = ET.fromstring(zf.read(rel_path))
    for rel in root.findall(f"{{{PPTX_REL_NS}}}Relationship"):
        rel_type = rel.attrib.get("Type", "")
        if not rel_type.endswith(PPTX_NOTES_REL_SUFFIX):
            continue
        target = rel.attrib.get("Target", "").strip()
        if not target:
            continue
        normalized = target.lstrip("/")
        if normalized.startswith("../"):
            normalized = normalized[3:]
        if not normalized.startswith("ppt/"):
            normalized = f"ppt/{normalized}"
        if normalized in zf.namelist():
            return normalized
    return None


def _parse_pptx_xml_fallback(path: Path, config: ParseConfig) -> tuple[list[_PptxSlideRecord], int]:
    records: list[_PptxSlideRecord] = []
    total_chars = 0
    with zipfile.ZipFile(path) as zf:
        for slide_number, slide_xml_path in _pptx_slide_xml_paths(zf):
            texts = _pptx_xml_text(zf.read(slide_xml_path))
            body_text = _normalize_text("\n\n".join(texts))
            notes_text = ""
            notes_path = _pptx_notes_xml_path(zf, slide_number)
            if notes_path:
                notes_text = _normalize_text("\n\n".join(_pptx_xml_text(zf.read(notes_path))))

            title = texts[0] if texts else f"Slide {slide_number}"
            record = _build_pptx_slide_record(
                slide_number=slide_number,
                slide_title=title,
                body_text=body_text,
                notes_text=notes_text,
                shape_count=0,
                min_chars_pptx_slide=config.min_chars_pptx_slide,
            )
            records.append(record)
            total_chars += len(record.text)
    return records, total_chars


def _pptx_records_to_docs(records: list[_PptxSlideRecord], base: dict) -> list[Document]:
    docs: list[Document] = []
    for record in records:
        metadata = dict(base)
        metadata["slide_number"] = record.slide_number
        metadata["slide_title"] = record.slide_title
        metadata["slide_has_notes"] = record.has_notes
        metadata["visual_enriched"] = False
        metadata["shape_count"] = record.shape_count
        docs.append(Document(page_content=record.text, metadata=metadata))
    return docs


def _ollama_url() -> str:
    return (os.getenv("OLLAMA_URL") or os.getenv("OLLAMA_HOST") or "http://127.0.0.1:11434").rstrip("/")


def _require_binary(name: str) -> None:
    if shutil.which(name) is None:
        raise RuntimeError(f"Required binary '{name}' not found")


def _render_pptx_slide_images(path: Path, temp_dir: Path) -> list[Path]:
    _require_binary("libreoffice")
    _require_binary("pdftoppm")

    pdf_dir = temp_dir / "pdf"
    img_dir = temp_dir / "slides"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    img_dir.mkdir(parents=True, exist_ok=True)

    proc = subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_dir),
            str(path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.strip() or "libreoffice conversion failed")

    pdf_path = pdf_dir / f"{path.stem}.pdf"
    if not pdf_path.exists():
        raise RuntimeError(f"Expected rendered PDF not found: {pdf_path}")

    prefix = img_dir / "slide"
    proc = subprocess.run(
        ["pdftoppm", "-png", "-r", "200", str(pdf_path), str(prefix)],
        capture_output=True,
        text=True,
        check=False,
    )
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.strip() or "pdftoppm failed")

    pattern = re.compile(r"^slide-(\d+)\.png$")
    slide_images = []
    for image_path in img_dir.glob("slide-*.png"):
        match = pattern.match(image_path.name)
        if not match:
            continue
        slide_images.append((int(match.group(1)), image_path))
    return [path for _, path in sorted(slide_images, key=lambda item: item[0])]


def _ollama_vision_summary(
    *,
    model: str,
    image_path: Path,
    slide_title: str,
    text_context: str,
) -> tuple[str, list[str]]:
    if requests is None:
        raise RuntimeError("requests is unavailable")

    prompt = (
        "Return JSON with keys visual_summary (string) and key_facts (array of short strings). "
        "Ground only in visible evidence from this slide image and provided text context. "
        f"Slide title: {slide_title}\n"
        f"Extracted text context:\n{text_context[:3500]}"
    )
    schema = {
        "type": "object",
        "properties": {
            "visual_summary": {"type": "string"},
            "key_facts": {"type": "array", "items": {"type": "string"}},
        },
        "required": ["visual_summary", "key_facts"],
    }
    with image_path.open("rb") as f:
        image_b64 = base64.b64encode(f.read()).decode("utf-8")

    payload = {
        "model": model,
        "prompt": prompt,
        "stream": False,
        "images": [image_b64],
        "format": schema,
    }
    response = requests.post(f"{_ollama_url()}/api/generate", json=payload, timeout=180)
    response.raise_for_status()
    body = response.json()
    data = json.loads(body.get("response", "{}"))
    summary = _normalize_text(str(data.get("visual_summary", "")))
    facts = []
    for fact in data.get("key_facts", []):
        fact_text = _normalize_text(str(fact))
        if fact_text:
            facts.append(fact_text)
    return summary, facts


def _pptx_should_try_vision(records: list[_PptxSlideRecord], config: ParseConfig) -> bool:
    if not records:
        return False
    if config.pptx_enable_vision:
        return True
    low_count = sum(1 for record in records if record.low_text)
    return (low_count / len(records)) >= config.pptx_vision_trigger_ratio


def _enrich_pptx_docs_with_vision(
    *,
    path: Path,
    docs: list[Document],
    records: list[_PptxSlideRecord],
    config: ParseConfig,
    issues: list[ParseIssue],
) -> bool:
    if config.pptx_vision_provider.lower() != "ollama":
        issues.append(
            ParseIssue(
                severity="warning",
                code="pptx_vision_skipped",
                message=f"Unsupported vision provider: {config.pptx_vision_provider}",
                suggested_action=PPTX_VISION_PROVIDER_GUIDANCE,
            )
        )
        return False

    low_records = [record for record in records if record.low_text]
    selected = low_records[: max(1, config.pptx_vision_max_slides)]
    if not selected:
        issues.append(
            ParseIssue(
                severity="warning",
                code="pptx_vision_skipped",
                message="No low-text slides selected for vision enrichment.",
                suggested_action="Enable --enable-vision to force slide enrichment.",
            )
        )
        return False

    try:
        with tempfile.TemporaryDirectory(prefix="myrag_pptx_vision_") as tmp:
            images = _render_pptx_slide_images(path, Path(tmp))
            image_by_slide = {idx: image for idx, image in enumerate(images, start=1)}

            enriched_count = 0
            for record in selected:
                image_path = image_by_slide.get(record.slide_number)
                if image_path is None:
                    continue
                summary, facts = _ollama_vision_summary(
                    model=config.pptx_vision_model,
                    image_path=image_path,
                    slide_title=record.slide_title,
                    text_context=record.text,
                )
                if not summary and not facts:
                    continue

                for doc in docs:
                    if doc.metadata.get("slide_number") != record.slide_number:
                        continue
                    chunks = [doc.page_content]
                    if summary:
                        chunks.append(f"Visual summary:\n{summary}")
                    if facts:
                        bullets = "\n".join(f"- {fact}" for fact in facts)
                        chunks.append(f"Visual key facts:\n{bullets}")
                    doc.page_content = _normalize_text("\n\n".join(chunks))
                    doc.metadata["visual_enriched"] = True
                    enriched_count += 1
                    break

            if enriched_count > 0:
                issues.append(
                    ParseIssue(
                        severity="warning",
                        code="pptx_vision_applied",
                        message=f"Applied vision enrichment to {enriched_count} low-text slide(s).",
                        suggested_action="Review enriched slide summaries for factual grounding.",
                    )
                )
                return True
    except Exception as exc:
        issues.append(
            ParseIssue(
                severity="warning",
                code="pptx_vision_failed",
                message=f"Vision enrichment failed: {exc}",
                suggested_action=f"{PPTX_LOCAL_VISION_GUIDANCE} {PPTX_VISION_PROVIDER_GUIDANCE}",
            )
        )
        return False

    issues.append(
        ParseIssue(
            severity="warning",
            code="pptx_vision_skipped",
            message="Vision enrichment ran but no additional visual content was extracted.",
            suggested_action="Retry with a stronger vision model or increase max vision slides.",
        )
    )
    return False


def parse_pptx_file(path: Path, config: ParseConfig) -> tuple[list[Document], FileParseResult]:
    extension = path.suffix.lower()
    base = _base_metadata(config, path)
    issues: list[ParseIssue] = []
    parser_used = "python-pptx"
    fallback_used = False

    primary_records: list[_PptxSlideRecord] = []
    primary_chars = 0
    primary_error: Exception | None = None
    try:
        primary_records, primary_chars = _parse_pptx_with_python_pptx(path, config)
    except Exception as exc:
        primary_error = exc
        text = str(exc).lower()
        missing_dep = "python-pptx" in text or "no module named" in text
        issues.append(
            ParseIssue(
                severity="warning",
                code="pptx_primary_parser_unavailable" if missing_dep else "pptx_primary_parser_failed",
                message=f"python-pptx parser failed: {exc}",
                suggested_action=(
                    "Install `python-pptx` for primary PPTX parsing. " + LOW_TEXT_GUIDANCE
                    if missing_dep
                    else LOW_TEXT_GUIDANCE
                ),
            )
        )

    chosen_records = primary_records
    chosen_chars = primary_chars
    if primary_error is not None or primary_chars < config.min_chars_pptx:
        fallback_used = True
        try:
            fallback_records, fallback_chars = _parse_pptx_xml_fallback(path, config)
            if fallback_chars >= chosen_chars:
                chosen_records = fallback_records
                chosen_chars = fallback_chars
                parser_used = "pptx_xml_fallback"
        except Exception as exc:
            issues.append(
                ParseIssue(
                    severity="warning",
                    code="pptx_xml_fallback_failed",
                    message=f"PPTX XML fallback failed: {exc}",
                    suggested_action="Validate the PPTX file and retry with repaired export.",
                )
            )

    docs = _pptx_records_to_docs(chosen_records, base)

    vision_applied = False
    if docs and _pptx_should_try_vision(chosen_records, config):
        vision_applied = _enrich_pptx_docs_with_vision(
            path=path,
            docs=docs,
            records=chosen_records,
            config=config,
            issues=issues,
        )

    status: ParseStatus = "success"
    if chosen_chars == 0 or len(docs) == 0:
        status = "failed"
        issues.append(
            ParseIssue(
                severity="error",
                code="pptx_no_extractable_text",
                message="No extractable text from PPTX file.",
                suggested_action=(
                    "Re-export the deck to PPTX/PDF, then retry. "
                    "For image-heavy decks, enable vision enrichment."
                ),
            )
        )
    else:
        low_count = sum(1 for record in chosen_records if record.low_text)
        low_ratio = (low_count / len(chosen_records)) if chosen_records else 0.0
        if chosen_chars < config.min_chars_pptx or (
            low_ratio >= config.pptx_vision_trigger_ratio and not vision_applied
        ):
            status = "warning"
            issues.append(
                ParseIssue(
                    severity="warning",
                    code="pptx_low_text",
                    message=(
                        "PPTX extraction quality is below threshold for part of the deck; "
                        "consider enabling selective vision enrichment."
                    ),
                    suggested_action=(
                        "Enable vision mode for low-text slides or enrich externally, then re-run."
                    ),
                )
            )

    _annotate_docs(docs, parser_used, fallback_used, status, chosen_chars)
    return docs, _result(
        path=path,
        extension=extension,
        status=status,
        parser_used=parser_used,
        fallback_used=fallback_used,
        extracted_chars=chosen_chars,
        documents_count=len(docs),
        issues=issues,
    )


def parse_knowledge_base(config: ParseConfig) -> tuple[list[Document], list[FileParseResult]]:
    if not config.knowledge_root.exists():
        raise FileNotFoundError(f"Knowledge root does not exist: {config.knowledge_root}")
    if not config.knowledge_root.is_dir():
        raise NotADirectoryError(f"Knowledge root is not a directory: {config.knowledge_root}")

    parser_map: dict[str, Callable[[Path, ParseConfig], tuple[list[Document], FileParseResult]]] = {
        ".pdf": parse_pdf_file,
        ".docx": parse_docx_file,
        ".xlsx": parse_xlsx_file,
        ".pptx": parse_pptx_file,
    }

    documents: list[Document] = []
    results: list[FileParseResult] = []

    for path in _iter_supported_files(config):
        ext = path.suffix.lower()
        parser = parser_map.get(ext)
        if parser is None:
            skipped = _result(
                path=path,
                extension=ext,
                status="skipped",
                parser_used="none",
                fallback_used=False,
                extracted_chars=0,
                documents_count=0,
                issues=[
                    ParseIssue(
                        severity="warning",
                        code="unsupported_extension",
                        message=f"Unsupported extension: {ext}",
                        suggested_action="Add a parser implementation for this file type.",
                    )
                ],
            )
            results.append(skipped)
            continue

        try:
            parsed_docs, parse_result = parser(path, config)
        except Exception as exc:
            parsed_docs = []
            parse_result = _result(
                path=path,
                extension=ext,
                status="failed",
                parser_used=parser.__name__,
                fallback_used=False,
                extracted_chars=0,
                documents_count=0,
                issues=[
                    ParseIssue(
                        severity="error",
                        code="unexpected_parse_error",
                        message=f"Unhandled parser error: {exc}",
                        suggested_action="Inspect parser logs and retry with an alternate parser.",
                    )
                ],
            )

        documents.extend(parsed_docs)
        results.append(parse_result)

    return documents, results
