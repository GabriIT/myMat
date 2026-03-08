from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from myMAT_app.parser.export_chunks import _chunk_id
from myMAT_app.parser.parser_config import ParseConfig
from myMAT_app.parser.parsers import parse_pptx_file

try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:
    Presentation = None
    Inches = None


def _make_sample_pptx(path: Path) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx not available")

    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Metal Replacement Intro"
    slide1.placeholders[1].text = (
        "This slide compares metal replacement options for housings and connectors."
    )
    notes = slide1.notes_slide.notes_text_frame
    notes.text = "Presenter note: emphasize cost and weight benefits."
    table = slide1.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(6), Inches(1.5)).table
    table.cell(0, 0).text = "Material"
    table.cell(0, 1).text = "Benefit"
    table.cell(1, 0).text = "PA"
    table.cell(1, 1).text = "Weight reduction"

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Scope"
    textbox = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    textbox.text_frame.text = "Short"

    prs.save(str(path))


@unittest.skipIf(Presentation is None, "python-pptx not installed")
class PptxParserUnitTests(unittest.TestCase):
    def test_parse_pptx_populates_slide_metadata_and_notes(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            path = root / "Metal_Replacement" / "demo.pptx"
            _make_sample_pptx(path)

            config = ParseConfig(
                knowledge_root=root,
                min_chars_pptx=1,
                min_chars_pptx_slide=1,
            )
            docs, result = parse_pptx_file(path, config)

            self.assertGreaterEqual(len(docs), 2)
            self.assertEqual(".pptx", result.extension)
            self.assertEqual(str(path), result.source_path)
            self.assertEqual("success", result.status)

            first = docs[0]
            self.assertEqual(1, first.metadata["slide_number"])
            self.assertEqual("Metal Replacement Intro", first.metadata["slide_title"])
            self.assertTrue(first.metadata["slide_has_notes"])
            self.assertFalse(first.metadata["visual_enriched"])
            self.assertIn("Speaker notes", first.page_content)

    def test_low_text_threshold_can_trigger_warning(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            path = root / "Metal_Replacement" / "demo.pptx"
            _make_sample_pptx(path)

            config = ParseConfig(
                knowledge_root=root,
                min_chars_pptx=10000,
                min_chars_pptx_slide=500,
            )
            _, result = parse_pptx_file(path, config)

            self.assertIn(result.status, {"warning", "failed"})
            issue_codes = {issue.code for issue in result.issues}
            self.assertIn("pptx_low_text", issue_codes)

    def test_primary_failure_uses_xml_fallback(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            path = root / "Metal_Replacement" / "demo.pptx"
            _make_sample_pptx(path)

            config = ParseConfig(
                knowledge_root=root,
                min_chars_pptx=1,
                min_chars_pptx_slide=1,
            )

            with patch(
                "myMAT_app.parser.parsers._parse_pptx_with_python_pptx",
                side_effect=RuntimeError("boom"),
            ):
                docs, result = parse_pptx_file(path, config)

            self.assertGreater(len(docs), 0)
            self.assertTrue(result.fallback_used)
            self.assertEqual("pptx_xml_fallback", result.parser_used)
            issue_codes = {issue.code for issue in result.issues}
            self.assertIn("pptx_primary_parser_failed", issue_codes)

    def test_chunk_id_includes_slide_number(self) -> None:
        m1 = {
            "source": "/tmp/demo.pptx",
            "source_name": "demo.pptx",
            "slide_number": 1,
        }
        m2 = {
            "source": "/tmp/demo.pptx",
            "source_name": "demo.pptx",
            "slide_number": 2,
        }
        self.assertNotEqual(_chunk_id(m1, 0), _chunk_id(m2, 0))


if __name__ == "__main__":
    unittest.main()
